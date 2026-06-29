import os
import subprocess
import shutil
import tempfile
import glob
from pptx import Presentation


class PowerPointProcessor:
    def __init__(self):
        pass

    def extract_text(self, ppt_path):
        """
        Extracts structural text elements cleanly.
        Runs identically on both local desktop and cloud server environments.
        """
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"Presentation file not found at: {ppt_path}")

        prs = Presentation(ppt_path)
        extracted_text = ""

        for idx, slide in enumerate(prs.slides):
            extracted_text += f"\n--- Slide {idx + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    clean_text = " ".join(shape.text.split())
                    extracted_text += clean_text + "\n"

        return extracted_text

    def export_slide_images(self, ppt_path, output_dir):
        """
        Renders each slide of the presentation to a PNG image.

        On Windows, uses PowerPoint COM automation (win32com) if available,
        since it gives the most accurate rendering of the original deck.

        On Linux (Streamlit Cloud and most servers), uses LibreOffice in
        headless mode to do the actual rasterization. This requires
        LibreOffice to be installed in the environment — on Streamlit Cloud,
        add a `packages.txt` file at the repo root containing the line:
            libreoffice

        Raises a clear, loud error if no rendering backend is available,
        rather than silently producing blank placeholder images.
        """
        os.makedirs(output_dir, exist_ok=True)

        # --- Track 1: Windows PowerPoint COM automation ---
        try:
            import win32com.client

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt_app.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)

            total_slides = presentation.Slides.Count
            for idx in range(1, total_slides + 1):
                slide = presentation.Slides(idx)
                output_image_path = os.path.join(output_dir, f"slide_{idx}.png")
                slide.Export(output_image_path, "PNG")

            presentation.Close()
            ppt_app.Quit()
            return total_slides

        except (ImportError, ModuleNotFoundError):
            # win32com isn't available on this system (e.g. Linux) — fall through
            # to the LibreOffice-based path below.
            pass

        # --- Track 2: LibreOffice headless rendering (Linux / Streamlit Cloud) ---
        return self._export_slide_images_libreoffice(ppt_path, output_dir)

    def _export_slide_images_libreoffice(self, ppt_path, output_dir):
        """
        Uses LibreOffice in headless mode to convert each slide of the
        presentation into a PNG image.

        Strategy:
        1. Convert the whole .pptx to a single multi-page PDF via LibreOffice.
        2. Rasterize each page of that PDF into a separate PNG using pdftoppm
           (part of the poppler-utils package).

        This two-step approach is far more reliable across LibreOffice
        versions than asking it to export PNGs directly, since direct
        PNG export from `soffice` only reliably produces the *first* slide.
        """
        soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice_path:
            raise RuntimeError(
                "LibreOffice ('soffice') was not found on this system, so slide "
                "images cannot be rendered. On Streamlit Cloud, add a "
                "'packages.txt' file at your repo root containing the line "
                "'libreoffice' (and 'poppler-utils' for PDF-to-image conversion), "
                "then redeploy."
            )

        pdftoppm_path = shutil.which("pdftoppm")
        if not pdftoppm_path:
            raise RuntimeError(
                "'pdftoppm' (from the poppler-utils package) was not found on "
                "this system. On Streamlit Cloud, add 'poppler-utils' to your "
                "'packages.txt' file at the repo root, then redeploy."
            )

        with tempfile.TemporaryDirectory() as tmp_dir:
            # Step 1: Convert .pptx -> .pdf
            convert_cmd = [
                soffice_path,
                "--headless",
                "--norestore",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                ppt_path,
            ]
            result = subprocess.run(
                convert_cmd,
                capture_output=True,
                text=True,
                timeout=180,
            )

            pdf_candidates = glob.glob(os.path.join(tmp_dir, "*.pdf"))
            if result.returncode != 0 or not pdf_candidates:
                raise RuntimeError(
                    f"LibreOffice failed to convert '{ppt_path}' to PDF. "
                    f"stdout: {result.stdout.strip()} | stderr: {result.stderr.strip()}"
                )

            pdf_path = pdf_candidates[0]

            # Step 2: Convert each PDF page -> PNG
            # pdftoppm writes files like: slide-1.png, slide-2.png, ...
            output_prefix = os.path.join(tmp_dir, "slide")
            ppm_cmd = [
                pdftoppm_path,
                "-png",
                "-r", "150",  # DPI; raise for higher-resolution slide images
                pdf_path,
                output_prefix,
            ]
            ppm_result = subprocess.run(
                ppm_cmd,
                capture_output=True,
                text=True,
                timeout=180,
            )

            rendered_pages = sorted(
                glob.glob(f"{output_prefix}-*.png"),
                key=lambda p: int(os.path.splitext(p)[0].rsplit("-", 1)[-1]),
            )

            if ppm_result.returncode != 0 or not rendered_pages:
                raise RuntimeError(
                    f"pdftoppm failed to rasterize '{pdf_path}'. "
                    f"stdout: {ppm_result.stdout.strip()} | stderr: {ppm_result.stderr.strip()}"
                )

            # Step 3: Move/rename rendered pages into the requested output_dir
            # with the filenames the rest of the app expects: slide_1.png, slide_2.png, ...
            total_slides = len(rendered_pages)
            for idx, page_path in enumerate(rendered_pages, start=1):
                dest_path = os.path.join(output_dir, f"slide_{idx}.png")
                shutil.copyfile(page_path, dest_path)

            return total_slides
